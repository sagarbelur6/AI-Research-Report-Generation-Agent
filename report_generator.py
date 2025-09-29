# report_generator.py
import os
from typing import Dict, Any, List
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from pptx.util import Inches, Pt
import textwrap
import datetime
import uuid

def save_pdf_report(report: Dict[str, Any], out_dir: str = "outputs") -> str:
    os.makedirs(out_dir, exist_ok=True)
    fname = f"report_{report['topic'].replace(' ','_')}_{int(report['generated_at'])}.pdf"
    path = os.path.join(out_dir, fname)
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    # Title
    c.setFont("Helvetica-Bold", 18)
    c.drawString(40, height - 60, f"Report: {report['topic']}")
    c.setFont("Helvetica", 10)
    c.drawString(40, height - 80, f"Generated: {datetime.datetime.utcfromtimestamp(report['generated_at']).isoformat()} UTC")
    # Executive summary + main text (wrap)
    text = report.get("revised_text") or report.get("text") or ""
    wrapper = textwrap.TextWrapper(width=100)
    lines = wrapper.wrap(text)
    y = height - 110
    c.setFont("Helvetica", 10)
    for line in lines:
        if y < 60:
            c.showPage()
            y = height - 60
            c.setFont("Helvetica", 10)
        c.drawString(40, y, line)
        y -= 12
    # visuals
    for vis in report.get("visuals", []):
        try:
            if y < 220:
                c.showPage()
                y = height - 60
            c.drawImage(ImageReader(vis), 60, y - 200, width=480, height=180)
            y -= 220
        except Exception:
            pass
    c.save()
    return path

def save_pptx(report: Dict[str, Any], out_dir: str = "outputs") -> str:
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation()
    # title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Report: {report['topic']}"
    subtitle.text = f"Generated at {datetime.datetime.utcfromtimestamp(report['generated_at']).isoformat()} UTC"

    # add summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    tx = slide.shapes.placeholders[1].text_frame
    tx.text = (report.get("revised_text") or report.get("text"))[:1000]

    # add visuals
    for vis in report.get("visuals", []):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        try:
            slide.shapes.add_picture(vis, left, top, width=Inches(8))
        except Exception:
            pass

    fname = f"slides_{report['topic'].replace(' ','_')}_{int(report['generated_at'])}.pptx"
    path = os.path.join(out_dir, fname)
    prs.save(path)
    return path

def save_visual(path: str, out_dir="outputs/visuals") -> str:
    os.makedirs(out_dir, exist_ok=True)
    import shutil
    basename = os.path.basename(path)
    dst = os.path.join(out_dir, basename)
    shutil.copyfile(path, dst)
    return dst
